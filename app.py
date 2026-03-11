import os
import zipfile

import pandas as pd
import seaborn as sns
import streamlit as st

# Optional exports
pdf_available = False
pptx_available = False
try:
    from fpdf import FPDF
    pdf_available = True
except ImportError:
    pdf_available = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    pptx_available = True
except ImportError:
    pptx_available = False

from eda import run_eda
from analysis import run_analysis
from report_generator import generate_excel

# Configure the web app title/icon (shown in the browser tab and the Streamlit header)
st.set_page_config(
    page_title="Automated Data Analyzer",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

st.title("📊 Automated Data Analyzer")

st.sidebar.markdown(
    """## Quick Start

- Upload a CSV or Excel file
- Pick the value/category columns
- Click **Generate Analysis**

> Tip: Try smaller datasets first to keep the UI responsive.
"""
)

# Helper functions

def ensure_output_dirs():
    os.makedirs("output/charts", exist_ok=True)


def save_eda_insights(df):
    """Save EDA stats and histograms to disk and return paths."""
    ensure_output_dirs()

    stats_path = "output/eda_summary.txt"
    with open(stats_path, "w", encoding="utf-8") as f:
        f.write("=== Column Data Types ===\n")
        f.write(df.dtypes.astype(str).to_string())
        f.write("\n\n=== Missing Values ===\n")
        f.write(df.isnull().sum().to_string())
        f.write("\n\n=== Statistical Summary ===\n")
        f.write(df.describe().to_string())

    chart_paths = []
    numeric_columns = df.select_dtypes(include=["int64", "float64"]).columns
    for col in numeric_columns:
        fig = sns.histplot(df[col].dropna(), kde=True)
        fig.figure.suptitle(f"Distribution of {col}")
        chart_path = f"output/charts/eda_{col}.png"
        fig.figure.savefig(chart_path)
        chart_paths.append(chart_path)
        # Clear the figure to free memory
        fig.figure.clf()

    return stats_path, chart_paths


def create_pdf_report(df, summary, analysis_chart_path, output_path="output/report.pdf"):
    if not pdf_available:
        raise RuntimeError("fpdf is not installed. Install with `pip install fpdf` to generate PDF reports.")

    ensure_output_dirs()
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)

    pdf.add_page()
    pdf.set_font("Arial", "B", 16)
    pdf.cell(0, 10, "Automated Data Analysis Report", ln=1)

    pdf.set_font("Arial", "", 12)
    pdf.cell(0, 8, "\nColumn data types:", ln=1)
    for col, dtype in df.dtypes.astype(str).items():
        pdf.cell(0, 6, f"- {col}: {dtype}", ln=1)

    pdf.ln(4)
    pdf.cell(0, 8, "Missing values (per column):", ln=1)
    for col, miss in df.isnull().sum().items():
        pdf.cell(0, 6, f"- {col}: {miss}", ln=1)

    pdf.ln(4)
    pdf.cell(0, 8, "Summary statistics (first rows):", ln=1)
    stats_lines = summary.head(10).to_string().splitlines()
    for line in stats_lines:
        pdf.cell(0, 5, line, ln=1)

    # Add analysis chart
    if os.path.exists(analysis_chart_path):
        pdf.add_page()
        pdf.set_font("Arial", "B", 14)
        pdf.cell(0, 10, "Analysis Visualization", ln=1)
        pdf.image(analysis_chart_path, x=15, w=180)

    pdf.output(output_path)
    return output_path


def create_pptx_report(df, summary, analysis_chart_path, output_path="output/report.pptx"):
    if not pptx_available:
        raise RuntimeError("python-pptx is not installed. Install with `pip install python-pptx` to generate PPTX reports.")

    ensure_output_dirs()
    prs = Presentation()

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Automated Data Analysis Report"

    # Summary slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "Data summary"
    body = slide.shapes.placeholders[1].text_frame

    body.text = "Column data types:"
    for col, dtype in df.dtypes.astype(str).items():
        p = body.add_paragraph()
        p.text = f"{col}: {dtype}"
        p.level = 1

    body.add_paragraph().text = "\nMissing values (per column):"
    for col, miss in df.isnull().sum().items():
        p = body.add_paragraph()
        p.text = f"{col}: {miss}"
        p.level = 1

    # Add chart slide
    if os.path.exists(analysis_chart_path):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = "Analysis Visualization"
        slide.shapes.add_picture(analysis_chart_path, Inches(1), Inches(1.5), width=Inches(8))

    prs.save(output_path)
    return output_path


def create_eda_package(report_path, analysis_chart_path, stats_path, chart_paths):
    ensure_output_dirs()

    zip_path = "output/eda_package.zip"
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED) as z:
        z.write(report_path, arcname=os.path.basename(report_path))
        z.write(analysis_chart_path, arcname=os.path.basename(analysis_chart_path))
        z.write(stats_path, arcname=os.path.basename(stats_path))
        for path in chart_paths:
            z.write(path, arcname=os.path.basename(path))

        # Add optional formats if available
        pdf_path = "output/report.pdf"
        if os.path.exists(pdf_path):
            z.write(pdf_path, arcname=os.path.basename(pdf_path))
        pptx_path = "output/report.pptx"
        if os.path.exists(pptx_path):
            z.write(pptx_path, arcname=os.path.basename(pptx_path))

    return zip_path


# Main UI

uploaded_file = st.file_uploader(
    "Upload CSV or Excel dataset",
    type=["csv","xlsx"]
)

if uploaded_file:

    # Persist the uploaded dataframe so users can iteratively change dtypes/feature engineering
    if (
        "uploaded_file_name" not in st.session_state
        or st.session_state.uploaded_file_name != uploaded_file.name
    ):
        if uploaded_file.name.endswith("csv"):
            st.session_state.df = pd.read_csv(uploaded_file)
        else:
            st.session_state.df = pd.read_excel(uploaded_file)
        st.session_state.uploaded_file_name = uploaded_file.name

    df = st.session_state.df

    st.subheader("Dataset Preview")
    st.dataframe(df.head())

    st.subheader("Exploratory Data Analysis")

    run_eda(df)

    # Allow adjusting column dtypes (e.g., string → datetime, numeric → category)
    st.subheader("Column Data Types")
    dtype_table = df.dtypes.astype(str).to_frame("dtype")
    st.dataframe(dtype_table)

    with st.form("dtype_form"):
        type_col = st.selectbox("Select column to change type", options=df.columns)
        target_type = st.selectbox(
            "Target type",
            options=["int", "float", "string", "category", "datetime"],
            index=2,
        )
        apply_dtype = st.form_submit_button("Apply data type change")

        if apply_dtype:
            try:
                if target_type == "datetime":
                    df[type_col] = pd.to_datetime(df[type_col], errors="coerce")
                else:
                    df[type_col] = df[type_col].astype(target_type)
                st.session_state.df = df
                st.success(f"Converted `{type_col}` to `{target_type}`")
                st.experimental_rerun()
            except Exception as e:
                st.error(f"Failed to convert `{type_col}` to `{target_type}`: {e}")

    # Feature engineering: optionally convert a numeric column into a categorical (binned) column
    with st.expander("Feature Engineering: convert numeric → categorical (binning)", expanded=False):
        numeric_cols = df.select_dtypes(include=["int64", "float64"]).columns
        if len(numeric_cols) == 0:
            st.info("No numeric columns available to bin.")
        else:
            fe_numeric_col = st.selectbox(
                "Select numeric column to bin",
                options=numeric_cols,
                help="Create a new categorical column by binning a numeric column."
            )

            binning_method = st.radio(
                "Binning method",
                options=["Equal width (pd.cut)", "Quantiles (pd.qcut)"],
                index=0,
                horizontal=True,
            )

            bin_count = st.slider("Number of bins", min_value=2, max_value=12, value=5)

            if st.button("Create binned column"):
                binned_col_name = f"{fe_numeric_col}_binned"
                try:
                    if "Quantiles" in binning_method:
                        df[binned_col_name] = pd.qcut(df[fe_numeric_col], q=bin_count, duplicates='drop')
                    else:
                        df[binned_col_name] = pd.cut(df[fe_numeric_col], bins=bin_count)

                    st.session_state.df = df
                    st.success(f"Created new column: `{binned_col_name}`")
                    st.dataframe(df[[fe_numeric_col, binned_col_name]].head())
                except Exception as e:
                    st.error(f"Failed to create binned column: {e}")

    st.subheader("Column Selection")

    numeric_columns = df.select_dtypes(include=['int64','float64']).columns
    categorical_columns = df.select_dtypes(include=['object','category']).columns

    value_column = st.selectbox("Value Column", numeric_columns)
    category_column = st.selectbox("Category Column", categorical_columns)

    if st.button("Generate Analysis"):

        summary, fig = run_analysis(
            df,
            category_column,
            value_column
        )

        st.subheader("Summary Table")
        st.dataframe(summary)

        st.subheader("Visualization")
        st.pyplot(fig)

        ensure_output_dirs()
        chart_path = "output/charts/chart.png"
        fig.savefig(chart_path)

        with open(chart_path, "rb") as file:
            st.download_button(
                "Download Chart",
                file,
                "chart.png"
            )

        generate_excel(df, summary)

        if pdf_available:
            try:
                pdf_path = create_pdf_report(df, summary, chart_path)
                with open(pdf_path, "rb") as file:
                    st.download_button(
                        "Download PDF Report",
                        file,
                        "report.pdf"
                    )
            except Exception as e:
                st.warning(f"PDF export failed: {e}")
        else:
            st.info("Install `fpdf` to enable PDF export.")

        if pptx_available:
            try:
                pptx_path = create_pptx_report(df, summary, chart_path)
                with open(pptx_path, "rb") as file:
                    st.download_button(
                        "Download PPTX Report",
                        file,
                        "report.pptx"
                    )
            except Exception as e:
                st.warning(f"PPTX export failed: {e}")
        else:
            st.info("Install `python-pptx` to enable PPTX export.")

        # EDA + visuals download bundle
        stats_path, eda_chart_paths = save_eda_insights(df)
        eda_zip = create_eda_package(
            report_path="output/report.xlsx",
            analysis_chart_path=chart_path,
            stats_path=stats_path,
            chart_paths=eda_chart_paths,
        )

        with open(eda_zip, "rb") as file:
            st.download_button(
                "Download EDA + visuals package",
                file,
                "eda_package.zip"
            )

        with open("output/report.xlsx","rb") as file:
            st.download_button(
                "Download Excel Report",
                file,
                "report.xlsx"
            )